from pptx import Presentation
import deepl
import tkinter as tk
from tkinter import filedialog, ttk, messagebox
import os
from pdf2docx import Converter
from PyPDF2 import PdfReader
import docx
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import tempfile
from openpyxl import Workbook
import pandas as pd

class DocumentTranslator:
    def __init__(self, auth_key):
        self.translator = deepl.Translator(auth_key)
        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
    def translate_text(self, text, source_lang, target_lang):
        """텍스트 번역 메서드"""
        try:
            result = self.translator.translate_text(
                text,
                source_lang=source_lang,
                target_lang=target_lang
            )
            return str(result)
        except Exception as e:
            print(f"번역 중 오류: {str(e)}")
            return text
            
    def perform_ocr(self, pdf_path, progress_callback=None):
        """PDF 파일에서 OCR 수행"""
        extracted_text = []
        
        try:
            poppler_path = r"C:\Program Files\Release-24.08.0-0\poppler-24.08.0\Library\bin"
            os.environ['PATH'] = poppler_path + os.pathsep + os.environ['PATH']
            
            images = convert_from_path(
                pdf_path,
                dpi=300
            )
            total_pages = len(images)
            
            for i, image in enumerate(images):
                try:
                    # 임시 파일 처리 개선
                    temp_dir = tempfile.gettempdir()
                    temp_file = os.path.join(temp_dir, f'ocr_temp_{i}.png')
                    
                    image.save(temp_file, 'PNG')
                    text = pytesseract.image_to_string(
                        Image.open(temp_file),
                        lang='kor+eng'
                    )
                    
                    if text.strip():
                        extracted_text.append(text)
                    
                    if progress_callback:
                        progress = ((i + 1) / total_pages) * 50
                        progress_callback(progress)
                    
                    # 임시 파일 삭제 시도
                    try:
                        os.remove(temp_file)
                    except:
                        pass
                        
                except Exception as e:
                    print(f"페이지 {i+1} OCR 오류: {str(e)}")
                    continue
                
            return '\n'.join(extracted_text)
            
        except Exception as e:
            print(f"OCR 처리 중 오류: {str(e)}")
            raise
        
    def translate_pdf(self, input_file, output_file, source_lang, target_lang, progress_callback=None):
        try:
            # 먼저 일반적인 PDF to DOCX 변환 시도
            temp_docx = input_file.replace('.pdf', '_temp.docx')
            cv = Converter(input_file)
            cv.convert(temp_docx)
            cv.close()
            
            # DOCX 파일에서 텍스트 추출 시도
            doc = docx.Document(temp_docx)
            total_paragraphs = len(doc.paragraphs)
            text_content = '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            
            # 텍스트가 충분하지 않으면 OCR 수행
            if len(text_content.split()) < 10:  # 단어 수가 10개 미만이면 OCR 실행
                print("텍스트 추출 실패, OCR 시작...")
                text_content = self.perform_ocr(input_file, progress_callback)
            
            # 새 DOCX 문서 생성
            new_doc = docx.Document()
            paragraphs = text_content.split('\n')
            total_paragraphs = len(paragraphs)
            processed = 0
            
            for paragraph in paragraphs:
                if paragraph.strip():
                    translated_text = self.translate_text(
                        paragraph,
                        source_lang,
                        target_lang
                    )
                    new_doc.add_paragraph(translated_text)
                    
                    processed += 1
                    if progress_callback:
                        progress = 50 + ((processed / total_paragraphs) * 50)  # 번역은 나머지 50%
                        progress_callback(progress)
            
            # 번역된 DOCX 저장
            new_doc.save(output_file)
            
            # 임시 파일 삭제
            if os.path.exists(temp_docx):
                os.remove(temp_docx)
                
            return processed
            
        except Exception as e:
            print(f"PDF 처리 중 오류: {str(e)}")
            raise

    def translate_pptx(self, input_file, output_file, source_lang, target_lang, progress_callback=None):
        prs = Presentation(input_file)
        total_texts = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            total_texts += 1
        
        processed = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    # 텍스트 프레임의 원래 속성 저장
                    original_width = shape.width
                    original_height = shape.height
                    original_wrap = text_frame.word_wrap
                    
                    for paragraph in text_frame.paragraphs:
                        if paragraph.text.strip():
                            try:
                                # 원본 서식 정보 저장
                                original_font = None
                                original_size = None
                                original_bold = None
                                original_italic = None
                                original_alignment = paragraph.alignment
                                original_spacing = paragraph.line_spacing
                                
                                if paragraph.runs:
                                    original_font = paragraph.runs[0].font.name
                                    original_size = paragraph.runs[0].font.size
                                    original_bold = paragraph.runs[0].font.bold
                                    original_italic = paragraph.runs[0].font.italic
                                
                                # 날짜 형식 확인 및 처리
                                if self.is_date(paragraph.text):
                                    translated_text = self.translate_date(
                                        paragraph.text, 
                                        source_lang, 
                                        target_lang
                                    )
                                else:
                                    translated_text = self.translate_text(
                                        paragraph.text,
                                        source_lang,
                                        target_lang
                                    )
                                
                                # 번역된 텍스트 적용하면서 서식 유지
                                if translated_text and translated_text != paragraph.text:
                                    paragraph.clear()  # 기존 runs 제거
                                    run = paragraph.add_run()
                                    run.text = translated_text
                                    
                                    # 서식 복원
                                    if original_font:
                                        run.font.name = original_font
                                    if original_size:
                                        run.font.size = original_size
                                    if original_bold is not None:
                                        run.font.bold = original_bold
                                    if original_italic is not None:
                                        run.font.italic = original_italic
                                    
                                    paragraph.alignment = original_alignment
                                    if original_spacing:
                                        paragraph.line_spacing = original_spacing
                                
                                processed += 1
                                if progress_callback:
                                    progress = (processed / total_texts) * 100
                                    progress_callback(progress)
                                    
                            except Exception as e:
                                print(f"번역 오류 (슬라이드 {slide.slide_id}): {str(e)}")
                                continue
                    
                    # 텍스트 프레임 크기 조정
                    text_frame.word_wrap = original_wrap
                    if shape.width != original_width or shape.height != original_height:
                        shape.width = original_width
                        shape.height = original_height
        
        prs.save(output_file)
        return processed

    def is_date(self, text):
        """날짜 형식인지 확인"""
        import re
        date_patterns = [
            r'\d{4}[-/.]\d{1,2}[-/.]\d{1,2}',  # YYYY-MM-DD
            r'\d{1,2}[-/.]\d{1,2}[-/.]\d{4}',  # DD-MM-YYYY
            r'\d{4}년\s*\d{1,2}월\s*\d{1,2}일',  # 한국어 날짜
        ]
        return any(re.search(pattern, text) for pattern in date_patterns)

    def translate_date(self, date_text, source_lang, target_lang):
        """날짜 형식 변환"""
        import re
        
        # 한국어 날짜를 표준 형식으로 변환
        if source_lang == "KO":
            match = re.search(r'(\d{4})년\s*(\d{1,2})월\s*(\d{1,2})일', date_text)
            if match:
                year, month, day = match.groups()
                if target_lang == "ES":
                    return f"{day}/{month}/{year}"
                elif target_lang in ["EN-US", "EN-GB"]:
                    return f"{month}/{day}/{year}"
                else:
                    return f"{year}-{month:02d}-{day:02d}"
        
        # 다른 형식의 날짜 처리
        date_match = re.search(r'(\d{1,4})[-/.](\d{1,2})[-/.](\d{1,4})', date_text)
        if date_match:
            parts = list(date_match.groups())
            # 연도가 뒤에 있는 경우 조정
            if len(parts[0]) == 2 or len(parts[0]) == 1:
                parts = parts[2::-1]
            
            if target_lang == "ES":
                return f"{parts[2]}/{parts[1]}/{parts[0]}"
            elif target_lang in ["EN-US", "EN-GB"]:
                return f"{parts[1]}/{parts[2]}/{parts[0]}"
            else:
                return f"{parts[0]}-{int(parts[1]):02d}-{int(parts[2]):02d}"
        
        return date_text

    def image_to_excel(self, image_path, output_path, progress_callback=None):
        """이미지에서 텍스트 추출하여 엑셀로 저장"""
        try:
            # 이미지 열기
            image = Image.open(image_path)
            
            # OCR 수행
            if progress_callback:
                progress_callback(30)  # 30% 진행
                
            text = pytesseract.image_to_string(
                image,
                lang='kor+eng'  # 한국어 + 영어 인식
            )
            
            if progress_callback:
                progress_callback(60)  # 60% 진행
            
            # 텍스트를 줄 단위로 분리
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            
            # 엑셀 워크북 생성
            wb = Workbook()
            ws = wb.active
            ws.title = "추출된 텍스트"
            
            # 각 줄을 엑셀에 입력
            for i, line in enumerate(lines, 1):
                ws.cell(row=i, column=1, value=line)
                
                if progress_callback:
                    progress = 60 + (i/len(lines)) * 40  # 나머지 40% 진행
                    progress_callback(progress)
            
            # 엑셀 파일 저장
            wb.save(output_path)
            
            return len(lines)
            
        except Exception as e:
            print(f"이미지 처리 중 오류: {str(e)}")
            raise

    def image_to_word(self, image_path, output_path, source_lang, target_lang, progress_callback=None):
        """이미지에서 텍스트 추출, 번역 후 Word로 저장"""
        try:
            # 이미지 열기
            image = Image.open(image_path)
            
            # OCR 수행
            if progress_callback:
                progress_callback(30)  # 30% 진행
                
            text = pytesseract.image_to_string(
                image,
                lang='kor+eng'  # 한국어 + 영어 인식
            )
            
            if progress_callback:
                progress_callback(50)  # 50% 진행
            
            # 텍스트를 줄 단위로 분리하고 번역
            lines = [line.strip() for line in text.split('\n') if line.strip()]
            translated_lines = []
            
            # Word 문서 생성
            doc = docx.Document()
            
            # 원본 이미지 추가
            doc.add_picture(image_path, width=docx.shared.Inches(6))
            doc.add_paragraph("\n") # 간격 추가
            
            # 각 줄을 번역하고 Word에 추가
            for i, line in enumerate(lines):
                if line.strip():
                    # 번역
                    translated_text = self.translate_text(line, source_lang, target_lang)
                    
                    # 원본과 번역문을 함께 저장
                    p = doc.add_paragraph()
                    p.add_run(f"원문: {line}\n").bold = True
                    p.add_run(f"번역: {translated_text}\n\n")
                    
                    if progress_callback:
                        progress = 50 + (i/len(lines)) * 50  # 나머지 50% 진행
                        progress_callback(progress)
            
            # Word 파일 저장
            doc.save(output_path)
            
            return len(lines)
                
        except Exception as e:
            print(f"이미지 처리 중 오류: {str(e)}")
            raise

class TranslatorGUI:
    def __init__(self):
        self.window = tk.Tk()
        self.window.title("문서 번역기")
        self.window.geometry("600x450")  # 창 크기 조정
        
        # DeepL이 지원하는 언어 코드로 수정
        self.languages = {
            "한국어": "KO",
            "영어(미국)": "EN-US",
            "영어(영국)": "EN-GB",
            "일본어": "JA",
            "중국어": "ZH",
            "스페인어": "ES",  
            "프랑스어": "FR",
            "독일어": "DE",
            "이탈리아어": "IT",
            "포르투갈어": "PT-PT",
            "포르투갈어(브라질)": "PT-BR",
            "러시아어": "RU",
            "네덜란드어": "NL",
            "폴란드어": "PL",
            "우크라이나어": "UK"
        }
        
        # 언어 선택 프레임
        self.lang_frame = tk.Frame(self.window)
        self.lang_frame.pack(pady=10)
        
        # 원본 언어 선택
        tk.Label(self.lang_frame, text="원본 언어:").grid(row=0, column=0, padx=5)
        self.source_lang = ttk.Combobox(
            self.lang_frame, 
            values=list(self.languages.keys()),
            width=15
        )
        self.source_lang.grid(row=0, column=1, padx=5)
        self.source_lang.set("한국어")
        
        # 목표 언어 선택
        tk.Label(self.lang_frame, text="목표 언어:").grid(row=0, column=2, padx=5)
        self.target_lang = ttk.Combobox(
            self.lang_frame, 
            values=list(self.languages.keys()),
            width=15
        )
        self.target_lang.grid(row=0, column=3, padx=5)
        self.target_lang.set("스페인어")
        
        # API 키 입력
        self.api_frame = tk.Frame(self.window)
        self.api_frame.pack(pady=10)
        
        tk.Label(self.api_frame, text="DeepL API Key:").pack(side=tk.LEFT)
        self.api_key = tk.Entry(self.api_frame, width=30)
        self.api_key.pack(side=tk.LEFT, padx=5)
        
        # 파일 선택 프레임
        self.file_frame = tk.Frame(self.window)
        self.file_frame.pack(pady=10)
        
        self.file_label = tk.Label(self.file_frame, text="선택된 파일: 없음")
        self.file_label.pack()
        
        self.select_button = tk.Button(
            self.file_frame,
            text="파일 선택",
            command=self.select_file
        )
        self.select_button.pack(pady=5)
        
        # 진행바
        self.progress = ttk.Progressbar(
            self.window,
            orient="horizontal",
            length=300,
            mode="determinate"
        )
        self.progress.pack(pady=20)
        
        # 번역 버튼
        self.translate_button = tk.Button(
            self.window,
            text="번역 시작",
            command=self.start_translation
        )
        self.translate_button.pack(pady=10)
        
        # 상태 표시
        self.status_label = tk.Label(self.window, text="")
        self.status_label.pack(pady=10)
        
        self.window.mainloop()
    
    def select_file(self):
        self.filename = filedialog.askopenfilename(
            filetypes=[
                ("지원되는 파일", "*.pptx;*.pdf;*.png;*.jpg;*.jpeg"),
                ("PowerPoint files", "*.pptx"),
                ("PDF files", "*.pdf"),
                ("Image files", "*.png;*.jpg;*.jpeg")
            ]
        )
        if self.filename:
            self.file_label.config(
                text=f"선택된 파일: {os.path.basename(self.filename)}"
            )
    
    def update_progress(self, value):
        self.progress["value"] = value
        self.window.update_idletasks()
    
    def start_translation(self):
        if not hasattr(self, 'filename'):
            messagebox.showerror("오류", "파일을 선택해주세요.")
            return
            
        api_key = self.api_key.get()
        if not api_key:
            messagebox.showerror("오류", "DeepL API 키를 입력해주세요.")
            return
        
        source_lang_code = self.languages[self.source_lang.get()]
        target_lang_code = self.languages[self.target_lang.get()]
        
        try:
            translator = DocumentTranslator(api_key)
            file_ext = os.path.splitext(self.filename)[1].lower()
            
            if file_ext in ['.png', '.jpg', '.jpeg']:
                output_file = self.filename.replace(file_ext, f'_translated_{target_lang_code}.docx')
                processed = translator.image_to_word(
                    self.filename,
                    output_file,
                    source_lang_code,
                    target_lang_code,
                    self.update_progress
                )
                success_message = f"이미지 번역 완료! {processed}줄의 텍스트가 번역되었습니다."
                
            elif file_ext == '.pdf':
                output_file = self.filename.replace('.pdf', f'_translated_{target_lang_code}.docx')
                processed = translator.translate_pdf(
                    self.filename,
                    output_file,
                    source_lang_code,
                    target_lang_code,
                    self.update_progress
                )
                success_message = f"PDF 번역 완료! {processed}개의 텍스트가 번역되었습니다."
                
            else:  # .pptx
                output_file = self.filename.replace('.pptx', f'_translated_{target_lang_code}.pptx')
                processed = translator.translate_pptx(
                    self.filename,
                    output_file,
                    source_lang_code,
                    target_lang_code,
                    self.update_progress
                )
                success_message = f"PPT 번역 완료! {processed}개의 텍스트가 번역되었습니다."
            
            self.status_label.config(text=success_message)
            messagebox.showinfo(
                "완료", 
                f"작업이 완료되었습니다.\n저장 위치: {output_file}"
            )
            
        except Exception as e:
            messagebox.showerror("오류", f"처리 중 오류 발생: {str(e)}")
        finally:
            self.translate_button.config(state='normal')
            self.progress["value"] = 0

if __name__ == "__main__":
    app = TranslatorGUI()
